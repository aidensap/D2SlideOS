import re
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime
from app.config import OUTPUT_DIR

# 中文字体（Windows 系统字体）
def _get_cn_font():
    available = {f.name for f in fm.fontManager.ttflist}
    for name in ["Microsoft YaHei", "SimHei", "SimSun", "Arial Unicode MS"]:
        if name in available:
            return name
    return None

_CN_FONT = _get_cn_font()

def _apply_cn_font(ax):
    if not _CN_FONT:
        return
    for item in [ax.title, ax.xaxis.label, ax.yaxis.label]:
        item.set_fontfamily(_CN_FONT)
    for label in ax.get_xticklabels() + ax.get_yticklabels():
        label.set_fontfamily(_CN_FONT)
    legend = ax.get_legend()
    if legend:
        for text in legend.get_texts():
            text.set_fontfamily(_CN_FONT)

DARK_BG = RGBColor(15, 23, 42)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(226, 232, 240)
ACCENT = RGBColor(56, 189, 248)


def _strip_markdown(text: str) -> str:
    text = re.sub(r'#{1,6}\s*', '', text)
    text = re.sub(r'\*{1,2}(.+?)\*{1,2}', r'\1', text)
    text = re.sub(r'`(.+?)`', r'\1', text)
    return text.strip()


def _add_title_slide(prs, report_label, insights_summary, lang="en"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG

    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = report_label
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    L = LABELS[lang]
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(3.4), Inches(8), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]
    generated_label = "生成时间" if lang == "zh" else "Generated"
    p2.text = f"{generated_label}：{datetime.now().strftime('%Y-%m-%d %H:%M')}  |  Powered by SAP AI Core"
    p2.font.size = Pt(13)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

    if insights_summary:
        tb3 = slide.shapes.add_textbox(Inches(1), Inches(4.4), Inches(8), Inches(1.8))
        tb3.text_frame.word_wrap = True
        p3 = tb3.text_frame.paragraphs[0]
        p3.text = f"{L['summary']}{_strip_markdown(insights_summary)[:200]}"
        p3.font.size = Pt(13)
        p3.font.color.rgb = ACCENT


LABELS = {
    "zh": {
        "sales_title": "月度区域销售额", "sales_x": "月份", "sales_y": "净销售额",
        "on_time": "准时", "late": "延迟", "delivery_title": "交货及时率",
        "summary": "AI 摘要：", "generated": "生成时间",
    },
    "en": {
        "sales_title": "Monthly Sales by Region", "sales_x": "Month", "sales_y": "Net Revenue",
        "on_time": "On Time", "late": "Late", "delivery_title": "Delivery Performance",
        "summary": "AI Summary: ", "generated": "Generated",
    },
}


def _render_chart(df, lang="en"):
    L = LABELS[lang]
    fig, ax = plt.subplots(figsize=(10, 5))
    fig.patch.set_facecolor("#0f172a")
    ax.set_facecolor("#1e293b")
    ax.tick_params(colors="#e2e8f0")
    for spine in ax.spines.values():
        spine.set_color("#334155")

    if lang == "zh" and _CN_FONT:
        matplotlib.rcParams["font.family"] = _CN_FONT

    try:
        cols = set(df.columns.tolist())

        if {"MONTH", "REGIO", "NETWR"}.issubset(cols):
            pivot = df.groupby(["MONTH", "REGIO"])["NETWR"].sum().unstack(fill_value=0)
            pivot.plot(kind="bar", ax=ax, colormap="tab10")
            ax.set_title(L["sales_title"], color="#e2e8f0", pad=10)
            ax.set_xlabel(L["sales_x"], color="#94a3b8")
            ax.set_ylabel(L["sales_y"], color="#94a3b8")
            ax.legend(facecolor="#1e293b", labelcolor="#e2e8f0")

        elif {"ON_TIME", "LATE", "MONTH"}.issubset(cols):
            x = range(len(df))
            ax.bar(x, df["ON_TIME"], label=L["on_time"], color="#34d399")
            ax.bar(x, df["LATE"], bottom=df["ON_TIME"], label=L["late"], color="#ef4444")
            ax.set_xticks(list(x))
            ax.set_xticklabels(df["MONTH"].tolist(), color="#e2e8f0")
            ax.set_title(L["delivery_title"], color="#e2e8f0", pad=10)
            ax.legend(facecolor="#1e293b", labelcolor="#e2e8f0")

        else:
            text_cols = df.select_dtypes(include="object").columns.tolist()
            num_cols = df.select_dtypes(include="number").columns.tolist()
            if text_cols and num_cols:
                label_col, val_col = text_cols[0], num_cols[0]
                plot_df = df.nlargest(10, val_col) if len(df) > 10 else df
                ax.barh(plot_df[label_col].astype(str), plot_df[val_col], color="#38bdf8")
                ax.set_title(f"{val_col} Top {len(plot_df)}", color="#e2e8f0", pad=10)
                ax.invert_yaxis()
                ax.tick_params(axis="y", labelsize=8, colors="#e2e8f0")
            else:
                raise ValueError("cannot identify chart type")

        tmp = f"{OUTPUT_DIR}/_tmp_chart.png"
        plt.tight_layout()
        plt.savefig(tmp, dpi=150, facecolor=fig.get_facecolor())
        plt.close()
        return tmp
    except Exception:
        plt.close()
        return None


def _add_chart_slide(prs, df, report_label, lang="en"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = report_label
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    img = _render_chart(df, lang)
    if img:
        slide.shapes.add_picture(img, Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
        os.remove(img)


LINES_PER_SLIDE = 14  # ~14 lines fit at Pt(14) + Pt(8) space in a 5.5" textbox


def _add_insights_slide(prs, lines, title, page=None, total=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title if (page is None or total == 1) else f"{title}  ({page}/{total})"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
    tb2.text_frame.word_wrap = True
    first = True
    for line in lines:
        para = tb2.text_frame.paragraphs[0] if first else tb2.text_frame.add_paragraph()
        first = False
        para.text = line
        para.font.size = Pt(14)
        para.font.color.rgb = LIGHT_GRAY
        para.space_after = Pt(8)


def _split_insights(insights):
    lines = []
    for line in insights.split("\n"):
        clean = _strip_markdown(line)
        if clean:
            lines.append(clean)
    chunks = []
    for i in range(0, max(len(lines), 1), LINES_PER_SLIDE):
        chunk = lines[i:i + LINES_PER_SLIDE]
        if chunk:
            chunks.append(chunk)
    return chunks or [[]]


def _add_screenshot_slide(prs, image_path, report_label, lang="en"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = report_label
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    slide.shapes.add_picture(image_path, Inches(0.2), Inches(0.9), Inches(9.6), Inches(5.8))


def generate_slides_from_screenshot(screenshot_path: str, report_label: str, insights: str, lang: str = "en") -> str:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7)

    first_line = next((l for l in insights.split("\n") if l.strip()), "")
    _add_title_slide(prs, report_label, first_line, lang)
    _add_screenshot_slide(prs, screenshot_path, report_label, lang)

    slide_title = "AI 数据洞察" if lang == "zh" else "AI Data Insights"
    chunks = _split_insights(insights)
    for i, chunk in enumerate(chunks, 1):
        _add_insights_slide(prs, chunk, slide_title, page=i, total=len(chunks))

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = f"{OUTPUT_DIR}/report_{timestamp}.pptx"
    prs.save(path)
    return path


def generate_slides(df, report_label, insights, lang="en"):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7)

    first_line = next((l for l in insights.split("\n") if l.strip()), "")
    _add_title_slide(prs, report_label, first_line, lang)
    _add_chart_slide(prs, df, report_label, lang)

    slide_title = "AI 数据洞察" if lang == "zh" else "AI Data Insights"
    chunks = _split_insights(insights)
    for i, chunk in enumerate(chunks, 1):
        _add_insights_slide(prs, chunk, slide_title, page=i, total=len(chunks))

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = f"{OUTPUT_DIR}/report_{timestamp}.pptx"
    prs.save(path)
    return path
